import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import *


# Written by Mehmet Enes Battal,
# with the intended use of automating the process of
# preparing an executive summary presentation.

CATEGORIES = ["security", "reliability", "performance-efficiency"]
FRAMEWORKS = ["NIST Cybersecurity Framework v1.1"]

csv_path = "test2.csv"
df = pd.read_csv(csv_path, dtype=object)

prs = Presentation()

# Set width and height to 16 and 9 inches.
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Slide Layouts
title_slide_layout = prs.slide_layouts[0]
blank_slide_layout = prs.slide_layouts[6]

# Slide 1
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
title.text = "Executive Summary"

# Add footer
slide.shapes.add_picture("footer.png", Inches(0), Inches(8.5), width=Inches(16))


# Drop the rows if relevant columns are null.
df.dropna(subset=["Cloud Provider", "Categories", "Risk Level"], inplace=True)

# Get the unique coud providers
providers = df["Cloud Provider"].unique().tolist()

fontSize = Pt(15)

# Slides for each cloud provider
for provider in providers:
    prv = df[df["Cloud Provider"] == provider]
    uniqueAccountCount = len(prv["Account ID"].unique())
    riskLevels = prv["Risk Level"].unique().tolist()

    left = top = Inches(0.5)
    width = height = Inches(1)

    slide = prs.slides.add_slide(blank_slide_layout)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    # Platform
    p = tf.paragraphs[0]
    p.text = "Platform: " + str(provider).upper()
    p.font.size = fontSize

    # Account number
    p = tf.add_paragraph()
    p.text = "Account Number: " + str(uniqueAccountCount)
    p.font.size = fontSize

    # Category
    p = tf.add_paragraph()
    p.text = "Category: " + ", ".join(CATEGORIES)
    p.font.size = fontSize

    # Framework
    p = tf.add_paragraph()
    p.text = "Frameworks: " + ", ".join(FRAMEWORKS)
    p.font.size = fontSize

    # Risk level
    p = tf.add_paragraph()
    p.text = "Risk Levels: " + ", ".join(riskLevels)
    p.font.size = fontSize

    print(tf.text)

    def add_data_labels(chart):
        plot = chart.plots[0]
        plot.has_data_labels = True
        for series in plot.series:
            values = series.values
            counter = 0
            for point in series.points:
                data_label = point.data_label
                data_label.has_text_frame = True
                data_label.text_frame.text = format(values[counter], ".0%")
                counter = counter + 1

    # Add chart for each category

    for idx, category in enumerate(CATEGORIES):
        tmp = prv[(prv['Categories'].str.contains(category))]

        chart_data = ChartData()
        chart_data.categories = ["Success", "Failure"]

        # Omitted not scored ones
        data = [len(tmp[tmp["Check Status"] == "SUCCESS"]), len(tmp[tmp["Check Status"] == "FAILURE"])]
        print(data)

        if sum(data) == 0:
            print("No checks for " + category + " in " + provider.upper())
            continue

        # Normalize
        norm = [float(i)/sum(data) for i in data]
        print(norm)

        chart_data.add_series("", norm)

        x, y = Inches(0.5 + 5 * idx), Inches(2)
        cx, cy = Inches(3), Inches(3.75)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data
        ).chart

        # Determine chart titles
        if "-" in category:
            l = [_.capitalize() for _ in category.split("-")]
            cTitle = "\n".join(l)
        else:
            cTitle = category.capitalize()

        chart.chart_title.text_frame.text = cTitle

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

        add_data_labels(chart)

        points = chart.plots[0].series[0].points
        fill = points[0].format.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 255, 0)

        fill = points[1].format.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 0, 0)

        
        # Info text
        txBox = slide.shapes.add_textbox(Inches(3.5 + 5 * idx), Inches(3), Inches(1), Inches(1))
        tf = txBox.text_frame
        
        p = tf.paragraphs[0]
        p.text = "Filter checked: " + str(sum(data))
        p.font.size = fontSize

        p = tf.add_paragraph()
        p.text = "Succeeded: " + str(data[0])
        p.font.size = fontSize

        p = tf.add_paragraph()
        p.text = "Failed: " + str(data[1])
        p.font.size = fontSize


    # Add footer
    slide.shapes.add_picture("footer.png", Inches(0), Inches(8.5), width=Inches(16))


output_path = "test.pptx"
prs.save(output_path)